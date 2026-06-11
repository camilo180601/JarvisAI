"""
document_creator.py — Crear y procesar documentos (cross-platform: Mac/Windows/Linux).

Crear:
  word        .docx (con markdown simple → títulos/listas)
  excel       .xlsx (varias hojas con headers/rows)
  powerpoint  .pptx (lista de slides con título + bullets)
  text        .txt

Procesar (lee PDF/DOCX/TXT y usa el cerebro IA configurado):
  summarize   resume un archivo o un texto
  translate   traduce un archivo o un texto (target=idioma)
  ocr         extrae texto de una imagen (Vision en Mac / tesseract si está)
"""
from __future__ import annotations
import os
from pathlib import Path
from datetime import datetime
from core.registry import tool


def _desktop() -> Path:
    """Escritorio cross-platform, con fallbacks."""
    for env in ("USERPROFILE", "HOME"):
        base = os.environ.get(env)
        if base:
            d = Path(base) / "Desktop"
            if d.exists():
                return d
    d = Path.home() / "Desktop"
    return d if d.exists() else Path.home()


def _safe(title: str) -> str:
    s = "".join(c for c in (title or "") if c.isalnum() or c == " ").strip().replace(" ", "_")
    return s or "Documento"


def _out_path(title: str, ext: str) -> Path:
    return _desktop() / f"{_safe(title)}_{datetime.now():%Y%m%d_%H%M%S}.{ext}"


def _read_text(path: str) -> str:
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"no existe {path}")
    suf = p.suffix.lower()
    if suf == ".pdf":
        from pypdf import PdfReader
        return "\n".join((pg.extract_text() or "") for pg in PdfReader(str(p)).pages)
    if suf in (".docx",):
        from docx import Document
        return "\n".join(par.text for par in Document(str(p)).paragraphs)
    return p.read_text(encoding="utf-8", errors="ignore")


@tool(
    name='document_creator',
    description="Crea y procesa documentos (cross-platform). Crear: word(.docx), excel(.xlsx), powerpoint(.pptx), text. Procesar archivos/textos con IA: summarize (resumir PDF/DOCX/TXT o texto), translate (traducir, target=idioma), ocr (texto desde imagen). Ej: 'hacé un Word con...', 'armá una presentación de...', 'resumí este PDF', 'traducí este documento al inglés'.",
    parameters={'type': 'OBJECT',
     'properties': {'action': {'type': 'STRING',
                               'description': 'word | excel | powerpoint | text | summarize | '
                                              'translate | ocr'},
                    'path': {'type': 'STRING',
                             'description': 'summarize/translate/ocr: ruta del archivo de entrada '
                                            '(PDF/DOCX/TXT/imagen)'},
                    'target': {'type': 'STRING',
                               'description': "translate: idioma destino (ej 'inglés', 'francés')"},
                    'save': {'type': 'BOOLEAN',
                             'description': 'summarize/translate: guardar el resultado en un .txt'},
                    'slides': {'type': 'ARRAY',
                               'description': 'powerpoint: lista de slides, cada uno con title '
                                              '(string) y bullets (array de strings).',
                               'items': {'type': 'OBJECT',
                                         'properties': {'title': {'type': 'STRING',
                                                                  'description': 'Título del slide'},
                                                        'bullets': {'type': 'ARRAY',
                                                                    'items': {'type': 'STRING'},
                                                                    'description': 'Viñetas del '
                                                                                   'slide'}}}},
                    'title': {'type': 'STRING',
                              'description': 'Title or filename of the document/spreadsheet'},
                    'content': {'type': 'STRING',
                                'description': 'For word / google_doc: full text content. Use ## '
                                               'Section for main headings, # SubSection for '
                                               'sub-headings, - item for bullet points, blank line '
                                               'between paragraphs.'},
                    'sheets': {'type': 'ARRAY',
                               'description': 'For excel / google_sheet: list of sheet objects. Each '
                                              'object has: name (string), headers (array of strings), '
                                              'rows (array of arrays with cell values).',
                               'items': {'type': 'OBJECT',
                                         'properties': {'name': {'type': 'STRING',
                                                                 'description': 'Sheet tab name'},
                                                        'headers': {'type': 'ARRAY',
                                                                    'items': {'type': 'STRING'},
                                                                    'description': 'Column headers'},
                                                        'rows': {'type': 'ARRAY',
                                                                 'items': {'type': 'ARRAY',
                                                                           'items': {'type': 'STRING'}},
                                                                 'description': 'Data rows'}}}},
                    'save_path': {'type': 'STRING',
                                  'description': 'Optional: full file path to save locally (e.g. '
                                                 'C:/Users/User/Desktop/report.docx). Defaults to '
                                                 '~/Documents/'}},
     'required': ['action', 'title']},
)
def document_creator(parameters: dict, player=None) -> str:
    action = (parameters.get("action") or "").lower()
    title = parameters.get("title", "Documento")
    content = parameters.get("content", "")

    try:
        # ───────── crear Word ─────────
        if action in ("word", "google_doc", "docx"):
            from docx import Document
            doc = Document()
            doc.add_heading(title, 0)
            for line in content.split("\n"):
                line = line.strip()
                if not line:
                    continue
                if line.startswith("## "):
                    doc.add_heading(line[3:], level=2)
                elif line.startswith("# "):
                    doc.add_heading(line[2:], level=1)
                elif line.startswith("- "):
                    doc.add_paragraph(line[2:], style="List Bullet")
                else:
                    doc.add_paragraph(line)
            out = _out_path(title, "docx")
            doc.save(out)
            return f"✓ Word creado: {out}"

        # ───────── crear Excel ─────────
        if action in ("excel", "google_sheet", "xlsx"):
            from openpyxl import Workbook
            sheets = parameters.get("sheets", [])
            if not sheets:
                return "Error: faltan 'sheets' (lista con name/headers/rows)."
            wb = Workbook()
            wb.remove(wb.active)
            for sd in sheets:
                ws = wb.create_sheet(title=(sd.get("name", "Hoja"))[:31])
                if sd.get("headers"):
                    ws.append(sd["headers"])
                for row in sd.get("rows", []):
                    ws.append(row)
            out = _out_path(title, "xlsx")
            wb.save(out)
            return f"✓ Excel creado: {out}"

        # ───────── crear PowerPoint ─────────
        if action in ("powerpoint", "pptx", "google_slides", "presentation"):
            from pptx import Presentation
            from pptx.util import Pt
            slides = parameters.get("slides", [])
            if not slides:
                return "Error: faltan 'slides' (lista con title/bullets)."
            prs = Presentation()
            # portada
            cover = prs.slides.add_slide(prs.slide_layouts[0])
            cover.shapes.title.text = title
            for sd in slides:
                s = prs.slides.add_slide(prs.slide_layouts[1])
                s.shapes.title.text = sd.get("title", "")
                body = s.placeholders[1].text_frame
                bullets = sd.get("bullets", [])
                body.text = bullets[0] if bullets else ""
                for b in bullets[1:]:
                    p = body.add_paragraph()
                    p.text = b
            out = _out_path(title, "pptx")
            prs.save(out)
            return f"✓ PowerPoint creado ({len(slides)} slides): {out}"

        # ───────── texto plano ─────────
        if action == "text":
            out = _out_path(title, "txt")
            out.write_text(f"{title}\n\n{content}", encoding="utf-8")
            return f"✓ Texto creado: {out}"

        # ───────── resumir / traducir ─────────
        if action in ("summarize", "summary", "resumir", "translate", "traducir"):
            src = parameters.get("path") or parameters.get("file")
            text = _read_text(src) if src else content
            if not text.strip():
                return "Error: dame 'path' (archivo) o 'content' (texto)."
            text = text[:30000]
            from core.llm_router import consult
            if action in ("translate", "traducir"):
                target = parameters.get("target") or parameters.get("lang") or "inglés"
                sysmsg = f"Sos un traductor profesional. Traducí al {target} de forma natural, sin agregar comentarios."
                prompt = text
            else:
                sysmsg = "Resumí en español de forma clara y concisa, con los puntos clave en viñetas."
                prompt = "Resumí esto:\n\n" + text
            if player:
                player.write_log("📄 Procesando documento con IA...")
            result, _ = consult(prompt, system=sysmsg, max_tokens=2000)
            if parameters.get("save"):
                out = _out_path(title or "resultado", "txt")
                out.write_text(result, encoding="utf-8")
                return f"✓ Guardado: {out}\n\n{result[:600]}"
            return result

        # ───────── OCR ─────────
        if action == "ocr":
            img = parameters.get("path") or parameters.get("image")
            if not img or not Path(img).exists():
                return "Error: dame 'path' a una imagen."
            try:
                import pytesseract
                from PIL import Image
                return pytesseract.image_to_string(Image.open(img)).strip() or "(sin texto detectado)"
            except Exception:
                return ("Para OCR instalá tesseract: brew install tesseract && "
                        ".venv/bin/pip install pytesseract")

        return ("Acción no reconocida. Crear: word, excel, powerpoint, text. "
                "Procesar: summarize, translate, ocr.")
    except Exception as e:
        return f"Error: {str(e)[:200]}"
